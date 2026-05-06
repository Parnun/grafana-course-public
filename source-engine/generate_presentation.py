#!/usr/bin/env python3
"""
Grafana for Developer Course — Presentation Generator
Theme: minimal, clean, white-blue  |  16:9  |  English
Run: python3 generate_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import os

# ─── Colour Palette ───────────────────────────────────────────────────────────
W    = RGBColor(0xFF, 0xFF, 0xFF)  # white
DB   = RGBColor(0x0D, 0x1B, 0x2A)  # dark navy (section / header bg)
PB   = RGBColor(0x00, 0x6F, 0xC6)  # primary blue
AB   = RGBColor(0x1A, 0x8C, 0xFF)  # accent / highlight blue
LB   = RGBColor(0xE6, 0xF2, 0xFF)  # light blue tint
OR   = RGBColor(0xF4, 0x6E, 0x00)  # Grafana orange
GR   = RGBColor(0x2D, 0xA4, 0x4E)  # green
RD   = RGBColor(0xD9, 0x2B, 0x2B)  # red
YL   = RGBColor(0xF5, 0xA6, 0x23)  # amber / pending
DG   = RGBColor(0x2D, 0x33, 0x3B)  # body text dark gray
MG   = RGBColor(0x6B, 0x72, 0x80)  # medium gray
SB   = RGBColor(0x00, 0x3D, 0x7D)  # strong blue
LC   = RGBColor(0xF0, 0xF4, 0xFF)  # near-white slide background

# ─── Presentation Setup ───────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]       # index 6 = blank layout

# ─── Low-level helpers ────────────────────────────────────────────────────────

def bg(slide, color):
    """Set solid background colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill=PB, line=None, lw=Pt(1)):
    """Add a rectangle (positions & sizes in Inches as floats)."""
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = lw
    else:
        shp.line.fill.background()
    return shp


def rrect(slide, l, t, w, h, fill=PB, line=None):
    """Add a rounded rectangle."""
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.5)
    else:
        shp.line.fill.background()
    return shp


def txt(slide, text, l, t, w, h,
        size=18, color=DG, bold=False, italic=False,
        align=PP_ALIGN.LEFT, wrap=True, name="Calibri"):
    """Add a text box."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = name
    return box


def stxt(shape, text, size=16, color=W, bold=False, align=PP_ALIGN.CENTER):
    """Set centred label text on an existing shape."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Calibri"


def notes(slide, text):
    """Add speaker notes."""
    slide.notes_slide.notes_text_frame.text = text


def arrow_right(slide, l, t, w=0.45, h=0.28, fill=MG):
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


# ─── Composite helpers ────────────────────────────────────────────────────────

def header_bar(slide, title, subtitle=None):
    """Dark top bar with title + optional subtitle."""
    rect(slide, 0, 0, 13.333, 1.3, fill=DB)
    rect(slide, 0, 1.3, 13.333, 0.055, fill=OR)   # orange accent line
    txt(slide, title, 0.45, 0.1, 11.5, 0.9,
        size=30, color=W, bold=True, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle, 0.45, 0.88, 10.0, 0.42,
            size=13, color=RGBColor(0xA0, 0xC8, 0xFF), align=PP_ALIGN.LEFT)


def bullets(slide, items, l=0.5, t=1.65, w=12.3, line_h=0.85, size=18, color=DG):
    """Render a list of bullet point strings."""
    marker = "  ›  "
    for i, item in enumerate(items):
        txt(slide, marker + item, l, t + i * line_h, w, line_h + 0.1,
            size=size, color=color)


def label_box(slide, text, l, t, w, h, fill=PB, tcolor=W, size=16, bold=False):
    s = rrect(slide, l, t, w, h, fill=fill)
    stxt(s, text, size=size, color=tcolor, bold=bold)
    return s


# ─── Slide Type Builders ──────────────────────────────────────────────────────

def cover_slide():
    sld = prs.slides.add_slide(BLANK)
    bg(sld, DB)
    rect(sld, 0, 5.4, 13.333, 2.1, fill=RGBColor(0x07, 0x12, 0x1E))
    rect(sld, 0, 2.75, 0.14, 4.5, fill=OR)   # vertical accent
    txt(sld, "Grafana for Developers",
        0.6, 1.0, 12.0, 1.5, size=54, color=W, bold=True)
    txt(sld, "A hands-on guide to observability, dashboards, and plugin development",
        0.65, 2.6, 10.5, 0.85, size=20, color=RGBColor(0xA0, 0xC8, 0xFF))
    txt(sld, "2-Day Training  ·  GitHub Codespaces  ·  Prometheus  ·  InfluxDB  ·  Grafana OSS",
        0.65, 5.7, 12.0, 0.6, size=13, color=MG)
    txt(sld, "◈  Grafana", 10.6, 6.6, 2.5, 0.65, size=22, color=OR, bold=True)
    notes(sld, "Welcome learners. This course covers the full Grafana developer journey — from container setup all the way to building custom plugins.")


def agenda_slide():
    sld = prs.slides.add_slide(BLANK)
    bg(sld, W)
    header_bar(sld, "Course Overview", "9 sections  ·  9 hands-on labs  ·  one shared Codespaces environment")
    sections = [
        ("1", "Docker Basics",                  "Containers, Compose, essential commands"),
        ("2", "Introduction to Grafana",        "Architecture, capabilities, first install"),
        ("3", "Working with Data Sources",      "Prometheus, InfluxDB, connection patterns"),
        ("4", "PromQL Basics",                  "Selectors, functions, aggregations"),
        ("5", "Queries and Panels",             "Explore → Panel → Dashboard"),
        ("6", "Advanced Dashboard Features",    "Variables, templating, drill-down"),
        ("7", "Working with Alerting",          "Rules, thresholds, state transitions"),
        ("8", "Grafana Plugins",                "Plugin types, scaffold, develop, live reload"),
        ("9", "Summary & Next Steps",           "Recap and open Q&A"),
    ]
    left_col  = sections[:5]
    right_col = sections[5:]
    for i, (num, sec, desc) in enumerate(left_col):
        label_box(sld, num, 0.4, 1.65 + i * 1.05, 0.44, 0.65, fill=PB, size=20, bold=True)
        txt(sld, sec,  1.0, 1.65 + i * 1.05, 5.2, 0.38, size=15, color=DB, bold=True)
        txt(sld, desc, 1.0, 1.97 + i * 1.05, 5.2, 0.35, size=12, color=MG)
    for i, (num, sec, desc) in enumerate(right_col):
        label_box(sld, num, 7.1, 1.65 + i * 1.05, 0.44, 0.65, fill=PB, size=20, bold=True)
        txt(sld, sec,  7.7, 1.65 + i * 1.05, 5.2, 0.38, size=15, color=DB, bold=True)
        txt(sld, desc, 7.7, 1.97 + i * 1.05, 5.2, 0.35, size=12, color=MG)
    notes(sld, "Day 1 covers sections 1–5; Day 2 covers sections 6–9.")


def section_divider(num, title, subtitle=None):
    sld = prs.slides.add_slide(BLANK)
    bg(sld, DB)
    rect(sld, 0, 0, 0.5, 7.5, fill=OR)       # orange left stripe
    txt(sld, f"Section {num}", 1.0, 1.8, 5.0, 0.65, size=18, color=MG)
    txt(sld, title, 1.0, 2.45, 11.5, 1.9, size=44, color=W, bold=True)
    if subtitle:
        txt(sld, subtitle, 1.0, 4.45, 10.5, 0.8, size=18,
            color=RGBColor(0xA0, 0xC8, 0xFF))
    notes(sld, f"Section {num}: {title}")
    return sld


def content_slide(title, subtitle, points, note=""):
    sld = prs.slides.add_slide(BLANK)
    bg(sld, W)
    header_bar(sld, title, subtitle)
    bullets(sld, points, t=1.65, line_h=1.0, size=19)
    if note:
        notes(sld, note)
    return sld


def content_slide_sm(title, subtitle, points, note=""):
    """Smaller text for 5-6 items."""
    sld = prs.slides.add_slide(BLANK)
    bg(sld, W)
    header_bar(sld, title, subtitle)
    bullets(sld, points, t=1.65, line_h=0.79, size=16)
    if note:
        notes(sld, note)
    return sld


def lab_slide(num, title, exercise_file, steps, outcomes, note=""):
    """Lab exercise slide."""
    sld = prs.slides.add_slide(BLANK)
    bg(sld, W)
    # Header
    rect(sld, 0, 0, 13.333, 1.3, fill=PB)
    rect(sld, 0, 1.3, 13.333, 0.055, fill=OR)
    txt(sld, f"Lab  ·  Exercise {num:02d}", 0.45, 0.06, 4.5, 0.46,
        size=13, color=RGBColor(0xA0, 0xC8, 0xFF))
    txt(sld, title, 0.45, 0.46, 12.0, 0.85, size=28, color=W, bold=True)
    # Steps column header
    rect(sld, 0.35, 1.5, 7.0, 0.42, fill=LB)
    txt(sld, "STEPS", 0.48, 1.53, 2.5, 0.35, size=11, color=PB, bold=True)
    for i, step in enumerate(steps):
        label_box(sld, str(i + 1), 0.38, 2.08 + i * 0.82,
                  0.44, 0.58, fill=PB, size=16, bold=True)
        txt(sld, step, 0.95, 2.1 + i * 0.82, 6.25, 0.65, size=14, color=DG)
    # Outcomes column header
    rect(sld, 7.55, 1.5, 5.4, 0.42, fill=RGBColor(0xE6, 0xFF, 0xF0))
    txt(sld, "EXPECTED OUTCOMES", 7.68, 1.53, 5.1, 0.35, size=11, color=GR, bold=True)
    for i, outcome in enumerate(outcomes):
        txt(sld, "✓  " + outcome, 7.6, 2.1 + i * 0.77, 5.2, 0.65, size=14, color=DG)
    # File ref footer
    txt(sld, "📄  " + exercise_file, 0.45, 6.88, 12.0, 0.5, size=11, color=MG)
    if note:
        notes(sld, note)
    return sld


# ─── Diagram builders ─────────────────────────────────────────────────────────

def diagram_docker_vs_vm(sld):
    """Side-by-side VM vs Container comparison."""
    txt(sld, "Virtual Machine", 0.5, 1.6, 5.5, 0.5,
        size=17, color=DB, bold=True, align=PP_ALIGN.CENTER)
    vm_layers = [
        ("App A",      RGBColor(0xF0, 0x60, 0x00)),
        ("Guest OS",   RGBColor(0xA0, 0x20, 0x20)),
        ("Hypervisor", RGBColor(0x60, 0x30, 0x10)),
        ("Hardware",   DG),
    ]
    for i, (label, c) in enumerate(vm_layers):
        s = rrect(sld, 0.9, 2.2 + i * 0.9, 4.7, 0.76, fill=c)
        stxt(s, label, size=15, color=W, bold=True)
    # Divider
    rect(sld, 6.55, 1.5, 0.06, 5.5, fill=RGBColor(0xCC, 0xCC, 0xCC))
    txt(sld, "vs", 6.35, 3.55, 0.7, 0.55, size=20, color=MG, bold=True, align=PP_ALIGN.CENTER)
    # Container side
    txt(sld, "Containers", 7.1, 1.6, 5.5, 0.5,
        size=17, color=DB, bold=True, align=PP_ALIGN.CENTER)
    for i, (label, c) in enumerate([
        ("App A", OR), ("App B", PB), ("App C", GR)
    ]):
        s = rrect(sld, 7.2 + i * 1.83, 2.2, 1.65, 1.55, fill=c)
        stxt(s, label, size=14, color=W, bold=True)
    s = rrect(sld, 7.2, 3.9, 5.35, 0.72, fill=PB)
    stxt(s, "Docker Engine", size=16, color=W, bold=True)
    s = rrect(sld, 7.2, 4.72, 5.35, 0.72, fill=DG)
    stxt(s, "Host OS  ·  Hardware", size=15, color=W, bold=True)
    # Insight bar
    rect(sld, 0.35, 5.7, 12.63, 0.65, fill=LB)
    txt(sld, "Containers share the host OS kernel — no Guest OS overhead, faster startup, lower resource use",
        0.55, 5.74, 12.1, 0.55, size=14, color=PB)


def diagram_docker_components(sld):
    """Image → Container → Volume → Network diagram."""
    items = [
        ("📦  Image",     "Blueprint / template\nRead-only layers",  PB,  0.35),
        ("🚀  Container", "Live running instance\nCreated from image", OR,  3.55),
        ("💾  Volume",    "Persistent data\nSurvives restarts",       GR,  6.75),
        ("🌐  Network",   "Service-to-service\ncommunication",        SB,  9.95),
    ]
    for label, desc, c, lx in items:
        s = rrect(sld, lx, 2.0, 2.9, 2.65, fill=c)
        stxt(s, label, size=15, color=W, bold=True)
        txt(sld, desc, lx + 0.15, 3.1, 2.6, 1.1, size=13, color=W, align=PP_ALIGN.CENTER)
    for x in [3.3, 6.5, 9.7]:
        arrow_right(sld, x, 3.17, 0.38, 0.28, fill=MG)
    rect(sld, 0.35, 5.0, 12.63, 0.55, fill=LB)
    txt(sld, "These four concepts appear in every Docker Compose file throughout the course.",
        0.55, 5.04, 12.1, 0.48, size=13, color=PB)


def diagram_compose_stack(sld):
    """Docker Compose multi-service stack diagram."""
    services = [
        ("Grafana  :3000",     PB,  0.55),
        ("Prometheus  :9090",  OR,  4.75),
        ("InfluxDB  :8086",    GR,  8.95),
    ]
    txt(sld, "docker-compose.yml", 0.4, 1.65, 12.5, 0.48,
        size=15, color=MG, align=PP_ALIGN.CENTER)
    for label, c, lx in services:
        s = rrect(sld, lx, 2.2, 3.5, 1.2, fill=c)
        stxt(s, label, size=16, color=W, bold=True)
    s = rect(sld, 0.35, 3.6, 12.63, 0.72, fill=RGBColor(0x0D, 0x1B, 0x2A))
    stxt(s, "Docker Internal Network  (compose_default)", size=14, color=W)
    rect(sld, 0.35, 4.6, 6.1, 0.75, fill=RGBColor(0x0D, 0x1F, 0x0D))
    txt(sld, "$ docker compose up -d", 0.55, 4.66, 5.8, 0.6,
        size=17, color=GR, bold=True)
    rect(sld, 7.1, 4.6, 6.0, 0.75, fill=RGBColor(0x1F, 0x0D, 0x0D))
    txt(sld, "$ docker compose down", 7.3, 4.66, 5.6, 0.6,
        size=17, color=RD, bold=True)
    rect(sld, 0.35, 5.7, 12.63, 0.55, fill=LB)
    txt(sld, "This same pattern is used in every lab exercise throughout the course.",
        0.55, 5.74, 12.1, 0.48, size=13, color=PB)


def diagram_grafana_arch(sld):
    """Grafana architecture flow: Browser → Grafana → Data Sources."""
    # Browser
    s = rrect(sld, 0.4, 2.8, 2.5, 1.2, fill=SB)
    stxt(s, "🌐  Browser\n(User)", size=15, color=W, bold=True)
    arrow_right(sld, 2.95, 3.22, 0.55, 0.3, fill=AB)
    # Grafana server
    s = rrect(sld, 3.55, 2.8, 3.0, 1.2, fill=OR)
    stxt(s, "◈  Grafana\nServer", size=17, color=W, bold=True)
    arrow_right(sld, 6.6, 3.22, 0.55, 0.3, fill=AB)
    # Data sources
    txt(sld, "Data Sources", 7.2, 1.75, 5.5, 0.5,
        size=15, color=DB, bold=True, align=PP_ALIGN.CENTER)
    for i, (label, c) in enumerate([
        ("Prometheus", OR), ("InfluxDB", PB), ("Loki / other", GR)
    ]):
        s = rrect(sld, 7.2, 2.25 + i * 1.15, 3.6, 0.9, fill=c)
        stxt(s, label, size=14, color=W, bold=True)
    # Features row
    for i, (label, c) in enumerate([
        ("Dashboards", PB), ("Explore", OR), ("Alerting", RD), ("Plugins", GR)
    ]):
        s = rrect(sld, 3.55 + i * 1.14, 4.55, 1.05, 0.7, fill=c)
        stxt(s, label, size=11, color=W, bold=True)
    rect(sld, 5.08, 4.0, 0.05, 0.56, fill=MG)
    rect(sld, 0.35, 5.55, 12.63, 0.58, fill=LB)
    txt(sld, "Grafana never stores metrics itself — it queries external systems and renders the results.",
        0.55, 5.59, 12.1, 0.48, size=13, color=PB)


def diagram_datasource_hub(sld):
    """Hub-and-spoke data source diagram."""
    # Center: Grafana
    s = rrect(sld, 5.2, 3.0, 2.9, 1.2, fill=OR)
    stxt(s, "◈  Grafana", size=20, color=W, bold=True)
    # Left data sources
    left_ds = [("Prometheus", OR, 1.7), ("InfluxDB", GR, 3.05), ("PostgreSQL", SB, 4.4)]
    for label, c, ty in left_ds:
        s = rrect(sld, 0.3, ty, 2.6, 0.85, fill=c)
        stxt(s, label, size=14, color=W, bold=True)
        arrow_right(sld, 2.95, ty + 0.22, 2.28, 0.28, fill=RGBColor(0xBB, 0xCC, 0xDD))
    # Right data sources
    right_ds = [("Loki", RD, 1.7), ("Elasticsearch", RGBColor(0xF0, 0xA0, 0x00), 3.05), ("Custom Plugin", MG, 4.4)]
    for label, c, ty in right_ds:
        s = rrect(sld, 10.4, ty, 2.6, 0.85, fill=c)
        stxt(s, label, size=14, color=W, bold=True)
        arrow_right(sld, 8.12, ty + 0.22, 2.28, 0.28, fill=RGBColor(0xBB, 0xCC, 0xDD))
    rect(sld, 0.35, 5.55, 12.63, 0.58, fill=LB)
    txt(sld, "All data sources follow the same Grafana pattern: pick target → enter credentials → test connection → query.",
        0.55, 5.59, 12.1, 0.48, size=13, color=PB)


def diagram_alert_states(sld):
    """Alert state machine: Normal → Pending → Firing."""
    states = [
        ("Normal",  "Query below\nthreshold",               GR, 1.0),
        ("Pending", "Threshold exceeded\n(within window)",  YL, 5.1),
        ("Firing",  "Threshold exceeded\n(window elapsed)", RD, 9.2),
    ]
    for label, desc, c, lx in states:
        s = rrect(sld, lx, 2.4, 3.0, 1.7, fill=c)
        stxt(s, label, size=20, color=W, bold=True)
        txt(sld, desc, lx + 0.2, 3.3, 2.6, 0.9, size=13, color=W, align=PP_ALIGN.CENTER)
    arrow_right(sld, 4.05, 3.05, 1.1, 0.4, fill=MG)
    arrow_right(sld, 8.15, 3.05, 1.1, 0.4, fill=MG)
    # Resolved path
    rect(sld, 1.0, 4.35, 11.2, 0.04, fill=MG)
    txt(sld, "Resolved (condition clears)  →  returns to Normal",
        1.0, 4.42, 11.2, 0.5, size=13, color=GR, italic=True)
    # Key concepts
    concepts = [
        ("Evaluation Interval\nHow often rule is checked", PB),
        ("Pending Period\nPrevents alert flapping",        SB),
        ("Silences\nSuppress notifications",               MG),
    ]
    for i, (label, c) in enumerate(concepts):
        s = rrect(sld, 0.4 + i * 4.2, 5.25, 3.8, 1.0, fill=c)
        stxt(s, label, size=13, color=W, bold=False)


def diagram_plugin_types(sld):
    """Three Grafana plugin type cards."""
    types = [
        ("Data Source Plugin",
         "Connect Grafana to\nany data backend.\nAdd new query editors.",
         PB,  0.4),
        ("Panel Plugin",
         "Create new\nvisualization types.\nCustom React components.",
         OR,  4.6),
        ("App Plugin",
         "Bundle data sources,\npanels & pages into\na complete product.",
         GR,  8.8),
    ]
    for label, desc, c, lx in types:
        s = rrect(sld, lx, 1.9, 3.8, 3.9, fill=c)
        stxt(s, label, size=18, color=W, bold=True)
        txt(sld, desc, lx + 0.2, 3.1, 3.4, 1.8, size=14, color=W, align=PP_ALIGN.CENTER)
    rect(sld, 0.35, 6.05, 12.63, 0.62, fill=LB)
    txt(sld, "Scaffold any type with:   npx @grafana/create-plugin@latest",
        0.55, 6.1, 12.1, 0.52, size=15, color=PB, bold=True)


def diagram_variables_flow(sld):
    """Variable dropdown → dynamic query → panels update."""
    txt(sld, "Variable Dropdowns", 0.4, 1.6, 12.5, 0.48,
        size=16, color=DB, bold=True, align=PP_ALIGN.CENTER)
    for i, (lbl, val) in enumerate([
        ("$job",      "prometheus"),
        ("$instance", "localhost:9090"),
        ("$interval", "5m"),
    ]):
        s = rrect(sld, 0.7 + i * 4.05, 2.18, 3.7, 0.85, fill=PB)
        stxt(s, f"{lbl}  =  {val}", size=15, color=W, bold=True)
    txt(sld, "↓", 6.35, 3.1, 0.9, 0.55, size=28, color=PB, bold=True, align=PP_ALIGN.CENTER)
    s = rect(sld, 0.7, 3.75, 11.9, 0.85, fill=DG)
    stxt(s, 'rate(http_requests_total{job="$job", instance="$instance"}[$interval])',
         size=16, color=GR)
    txt(sld, "↓", 6.35, 4.65, 0.9, 0.55, size=28, color=PB, bold=True, align=PP_ALIGN.CENTER)
    txt(sld, "All panels update simultaneously", 0.4, 5.22, 12.5, 0.45,
        size=14, color=DB, bold=True, align=PP_ALIGN.CENTER)
    for i, label in enumerate(["Panel A", "Panel B", "Panel C", "Panel D"]):
        s = rrect(sld, 0.6 + i * 3.1, 5.72, 2.8, 1.0, fill=LB, line=PB)
        stxt(s, label, size=14, color=PB, bold=True)


def diagram_drill_down(sld):
    """Dashboard drill-down / data link flow."""
    s = rrect(sld, 0.35, 1.9, 4.7, 2.6, fill=PB)
    stxt(s, "Ops Overview\n(Table Panel)\nAll Targets", size=15, color=W, bold=True)
    arrow_right(sld, 5.1, 2.98, 0.75, 0.4, fill=OR)
    txt(sld, "Click row\n(Data Link)", 5.08, 2.58, 1.3, 0.52, size=11, color=MG, align=PP_ALIGN.CENTER)
    s = rrect(sld, 5.9, 1.9, 4.7, 2.6, fill=OR)
    stxt(s, "Target Detail\n(Graphs)\n$instance", size=15, color=W, bold=True)
    # Dashboard link label
    s = rrect(sld, 10.75, 2.2, 2.2, 0.65, fill=LB, line=PB)
    stxt(s, "Dashboard Link →", size=12, color=PB, bold=True)
    s = rrect(sld, 10.75, 3.05, 2.2, 0.65, fill=LB, line=PB)
    stxt(s, "Time range passed", size=11, color=DG)
    # Context row
    txt(sld, "Context preserved across navigation:", 0.35, 4.75, 12.5, 0.45,
        size=14, color=DB, bold=True)
    for i, ctx in enumerate(["Time Range", "Variable Values", "Selected Data Point"]):
        s = rrect(sld, 0.4 + i * 4.15, 5.3, 3.8, 0.9, fill=LB, line=PB)
        stxt(s, ctx, size=14, color=PB, bold=True)


# ─── Build all slides ─────────────────────────────────────────────────────────

def main():

    # ── Cover + Agenda ────────────────────────────────────────────────────────
    cover_slide()
    agenda_slide()

    # ══ SECTION 1: Docker Basics ══════════════════════════════════════════════
    section_divider(1, "Docker Basics", "Containers  ·  Images  ·  Compose")

    # 1.1 Docker vs VM (diagram)
    sld = prs.slides.add_slide(BLANK)
    bg(sld, W)
    header_bar(sld, "What is Docker and Containers?", "Section 1  ·  Docker Basics")
    diagram_docker_vs_vm(sld)
    notes(sld, "Open with the big picture — containers give every learner the same baseline without machine-specific installs.")

    # 1.2 How Containers Work
    content_slide(
        "How Containers Work",
        "Processes, not virtual machines",
        [
            "Containers share the host OS kernel — no Guest OS overhead",
            "Each container has its own isolated process, network, and filesystem",
            "Many containers run simultaneously on a single machine",
            "Much faster to start and stop than full virtual machines",
        ],
        "Correct the VM misconception. A container is an isolated process managed by the Docker engine.",
    )

    # 1.3 Key Components (diagram)
    sld = prs.slides.add_slide(BLANK)
    bg(sld, W)
    header_bar(sld, "Key Components", "Image  ·  Container  ·  Volume  ·  Network")
    diagram_docker_components(sld)
    notes(sld, "Establish shared vocabulary. Learners will see these terms in every Docker Compose file.")

    # 1.4 Docker Desktop
    content_slide(
        "Docker Desktop",
        "Local GUI — same engine, different surface",
        [
            "GUI for managing containers, images, and volumes on your laptop",
            "Uses the same Docker engine as GitHub Codespaces",
            "This course runs in Codespaces so everyone starts from the same baseline",
            "Docker Desktop is what learners will use after the course for local work",
        ],
        "Codespaces reduces risk from machine differences and makes instructor support easier.",
    )

    # 1.5 Essential Commands
    content_slide_sm(
        "Essential Docker Commands",
        "Your debugging survival kit",
        [
            "docker ps                       — list running containers",
            "docker images                   — list locally available images",
            "docker logs <name>              — view container log output",
            "docker exec -it <name> sh       — open an interactive shell",
            "docker inspect <name>           — view container metadata",
        ],
        "Goal: give learners enough to check status and debug on their own during labs.",
    )

    # 1.6 Docker Compose (diagram)
    sld = prs.slides.add_slide(BLANK)
    bg(sld, W)
    header_bar(sld, "Docker Compose for Multi-Service Stacks",
               "One file  ·  all services  ·  one command")
    diagram_compose_stack(sld)
    notes(sld, "Direct bridge to the Grafana labs. Almost every exercise uses Compose to bring up Grafana, Prometheus, and InfluxDB together.")

    # ══ SECTION 2: Introduction to Grafana ═══════════════════════════════════
    section_divider(2, "Introduction to Grafana",
                    "Explore  ·  Dashboards  ·  Alerting  ·  Plugins")

    # 2.1 What is Grafana?
    content_slide(
        "What is Grafana?",
        "Visualize. Explore. Alert.",
        [
            "Open-source observability and data visualization platform",
            "Connects to many data source types: Prometheus, InfluxDB, Loki, and more",
            "Central place for querying, visualizing, and answering questions from data",
            "Used by developers, SREs, and platform teams worldwide",
        ],
        "Introduce from a business perspective — Grafana is not just a chart tool but a central observability hub.",
    )

    # 2.2 Core Capabilities (4-pillar card)
    sld = prs.slides.add_slide(BLANK)
    bg(sld, W)
    header_bar(sld, "Core Capabilities of Grafana", "The four pillars of this course")
    pillars = [
        ("🔍  Explore",    "Experiment with\nqueries, view live data",       PB,  0.35),
        ("📊  Dashboards", "Organize data for\nfast decision-making",         OR,  3.6),
        ("🔔  Alerting",   "Detect important\nconditions automatically",      RD,  6.85),
        ("🔌  Plugins",    "Extend Grafana with\ncustom capabilities",        GR, 10.1),
    ]
    for label, desc, c, lx in pillars:
        s = rrect(sld, lx, 1.9, 2.9, 3.5, fill=c)
        stxt(s, label, size=17, color=W, bold=True)
        txt(sld, desc, lx + 0.15, 3.15, 2.6, 1.5, size=14, color=W, align=PP_ALIGN.CENTER)
    notes(sld, "Connect these four pillars to what learners will actually do. This builds a mental map from the start.")

    # 2.3 Architecture (diagram)
    sld = prs.slides.add_slide(BLANK)
    bg(sld, W)
    header_bar(sld, "Basic Grafana Architecture", "Browser  →  Grafana  →  Data Sources")
    diagram_grafana_arch(sld)
    notes(sld, "Make data flow clear before first install. Grafana does not store metrics — it queries external systems.")

    # ══ SECTION 3: Installation & Configuration ═══════════════════════════════
    section_divider(3, "Installation & Configuration",
                    "GitHub Codespaces  ·  docker compose up  ·  Port 3000")

    # 3.1 Preparing the Environment
    content_slide(
        "Preparing the Environment",
        "GitHub Codespaces as the shared lab platform",
        [
            "GitHub Codespaces provides an identical environment for every learner",
            "Docker and Docker Compose are pre-installed — no local setup required",
            "Port forwarding automatically exposes services like Grafana on port 3000",
            "Key files in every lab:   .env   and   docker-compose.yml",
        ],
        "Codespaces reduces risk from machine differences and simplifies instructor support.",
    )

    # 3.2 Running Grafana
    content_slide_sm(
        "Installing and Running Grafana",
        "Start  →  Check  →  Open  →  Log In",
        [
            "cd  labfiles/01-setup-grafana-codespaces/",
            "docker compose up -d           (start all services in the background)",
            "docker compose ps               (verify all containers are Running)",
            "Open port 3000 in the Codespaces Ports panel",
            "Log in with:  admin / admin   and change password on first login",
        ],
        "Reinforce this same start-up pattern — learners will repeat it in every exercise.",
    )

    # 3.3 Lab 01
    lab_slide(
        1,
        "Setup Grafana in Codespaces",
        "Exercises/01-setup-grafana-codespaces.md",
        [
            "Clone the repo and open it in GitHub Codespaces",
            "Run docker compose up -d in labfiles/01-setup-grafana-codespaces/",
            "Forward port 3000 and open Grafana in the browser",
            "Inspect the Docker volume to understand data persistence",
        ],
        [
            "Grafana running and accessible on port 3000",
            "Login with admin credentials confirmed",
            "Port forwarding understood",
            "Persistent Docker volume observed",
        ],
        "Allow everyone to reach a working Grafana before moving to data sources.",
    )

    # ══ SECTION 4: Working with Data Sources ══════════════════════════════════
    section_divider(4, "Working with Data Sources",
                    "Prometheus  ·  InfluxDB  ·  Connection Patterns")

    # 4.1 What is a Data Source? (diagram)
    sld = prs.slides.add_slide(BLANK)
    bg(sld, W)
    header_bar(sld, "What is a Data Source?",
               "The connection point between Grafana and your data")
    diagram_datasource_hub(sld)
    notes(sld, "Graphs in Grafana always start by picking the right data source, not by choosing a panel type.")

    # 4.2 Prometheus
    content_slide(
        "Prometheus Data Source",
        "Time-series metrics  ·  PromQL",
        [
            "Prometheus scrapes and stores time-series metrics from configured targets",
            "Query language: PromQL — functional, built for time-series aggregation",
            "In Docker Compose, the URL is:   http://prometheus:9090",
            "Explore → Code mode → type PromQL to see live metric data",
        ],
        "Inside Compose, use the service name 'prometheus' instead of localhost.",
    )

    # 4.3 Lab 02
    lab_slide(
        2,
        "Prometheus Data Source",
        "Exercises/02-prometheus-data-source.md",
        [
            "Start the lab stack with docker compose up -d",
            "Grafana → Connections → Add new data source → Prometheus",
            "Set URL to  http://prometheus:9090  →  Save & test",
            "Open Explore and run:  up   and   prometheus_build_info",
        ],
        [
            "Prometheus data source added",
            "Connection test passed (green)",
            "First PromQL query result visible in Explore",
        ],
        "First data source learners connect themselves. Ensure everyone uses the correct URL and sees results in Explore.",
    )

    # 4.4 InfluxDB + Flux
    content_slide(
        "InfluxDB Data Source with Flux",
        "InfluxDB 2.x  ·  Flux query language",
        [
            "InfluxDB 2.x uses Flux — a functional scripting language for data pipelines",
            "Required connection fields: URL, Organization, Default Bucket, and Token",
            "Select 'Flux' as the query language when configuring the data source",
            "Compares a different paradigm alongside PromQL in the same Grafana UI",
        ],
        "Same Grafana connection pattern — only the language for talking to the backend changes.",
    )

    # 4.5 Lab 03
    lab_slide(
        3,
        "InfluxDB + Flux Data Source",
        "Exercises/03-influxdb-flux-data-source.md",
        [
            "Run the seed script:   ./scripts/seed-influx.sh",
            "Grafana → Connections → Add InfluxDB data source (Flux mode)",
            "Fill in URL, Org, Token, and Default Bucket",
            "Query the  lab_metrics  measurement in Explore",
        ],
        [
            "InfluxDB data seeded successfully",
            "InfluxDB data source connected (Flux mode)",
            "Basic Flux query returns results in Explore",
        ],
        "Show the query language difference while keeping the same Grafana workflow.",
    )

    # ══ SECTION 5: PromQL Basics ══════════════════════════════════════════════
    section_divider(5, "PromQL Basics",
                    "Selectors  ·  Functions  ·  Aggregations")

    # 5.1 What is PromQL?
    content_slide(
        "What is PromQL?",
        "The query language built into Prometheus",
        [
            "Functional language for selecting and aggregating time-series data",
            "Every expression returns an instant vector, range vector, or scalar",
            "Used in Grafana Explore, panel editors, and alert rule definitions",
            "Not SQL — designed around labelled metric streams and time windows",
        ],
        "PromQL is not a SQL dialect. Built around the idea of time-series data streams.",
    )

    # 5.2 Selectors
    content_slide_sm(
        "PromQL: Basic Metric Selectors",
        "Select, filter, narrow",
        [
            "up                                       — 1 = target reachable, 0 = down",
            '{job="prometheus"}                 — label selector to filter results',
            'up{job="prometheus",instance="…"}  — combine multiple selectors',
            "prometheus_build_info              — build metadata per Prometheus instance",
            "[5m]  range selector             — look back 5 minutes of samples",
        ],
        "Start with 'up' — every instance exposes it. Show how selectors narrow results.",
    )

    # 5.3 Functions & Aggregations
    content_slide_sm(
        "PromQL: Useful Functions & Aggregations",
        "Four patterns that cover most real-world queries",
        [
            "rate(metric[5m])              — per-second rate; use for counter metrics",
            "increase(metric[1m])          — total counter change; useful in alert thresholds",
            "sum by (label) (metric)       — aggregate across instances grouped by label",
            "avg_over_time(metric[5m])     — average gauge value over a rolling window",
            "Rule: never plot raw counter values — always use rate() or increase()",
        ],
        "These four patterns cover the majority of real-world PromQL queries.",
    )

    # 5.4 PromQL in Explore
    content_slide(
        "PromQL in Grafana Explore",
        "Builder mode vs Code mode",
        [
            "Select Prometheus as the data source in the Explore page",
            "Builder mode — visual form for discovery and learning",
            "Code mode — type raw PromQL for complex expressions",
            "Time range picker to zoom in; Query Inspector to view raw JSON",
        ],
        "Demonstrate switching modes. By the end, learners should feel confident reading and writing raw PromQL.",
    )

    # ══ SECTION 6: Queries and Panels ════════════════════════════════════════
    section_divider(6, "Queries and Panels", "Explore  →  Panel  →  Dashboard")

    # 6.1 Queries in Grafana
    content_slide(
        "Queries in Grafana",
        "Start with the question, not the chart type",
        [
            "Begin in Explore — experiment freely before committing to a panel",
            "Choose the data source first, then write a query that answers a question",
            "A good query in Explore transfers directly to a dashboard panel",
            "The query is the foundation — visualization follows from the data",
        ],
        "Shift mindset: 'what question am I answering?' not 'what graph do I want?'",
    )

    # 6.2 Designing Panels
    content_slide(
        "Designing and Building Panels",
        "Match the visualization to the question",
        [
            "Time Series — trends and changes over a time range",
            "Stat — a single current value read at a glance",
            "Table — comparison across multiple dimensions or instances",
            "Panel title, legend labels, and layout matter as much as the query",
        ],
        "Visualization choice is design, not aesthetics. Match what the viewer needs to decide.",
    )

    # 6.3 Lab 04
    lab_slide(
        4,
        "Queries and Panels",
        "Exercises/04-queries-and-panels.md",
        [
            "Import the starter dashboard (dashboards/starter.json)",
            "Add a Time Series panel:  up{job='prometheus'}",
            "Add a Stat panel:  prometheus_build_info",
            "Set panel titles, configure legends, and adjust layout",
        ],
        [
            "First real dashboard built",
            "Time Series vs Stat panel differences understood",
            "PromQL queries rendered on dashboard panels",
        ],
        "Allow extra time here. Learners start seeing how queries become real, usable dashboards.",
    )

    # ══ SECTION 7: Advanced Dashboard Features ════════════════════════════════
    section_divider(7, "Advanced Dashboard Features",
                    "Variables  ·  Templating  ·  Drill-Down Links")

    # 7.1 Variables Concepts (diagram)
    sld = prs.slides.add_slide(BLANK)
    bg(sld, W)
    header_bar(sld, "Dashboard Variables",
               "Eliminate hard-coded values  ·  Enable dynamic filtering")
    diagram_variables_flow(sld)
    notes(sld, "Variables let one dashboard cover many services, preventing dashboard sprawl.")

    # 7.2 Lab 05
    lab_slide(
        5,
        "Dashboard Variables with InfluxDB",
        "Exercises/05-dashboard-variables.md",
        [
            "Seed multi-region InfluxDB data with ./scripts/seed-influx.sh",
            "Create a 'service' Query variable using a Flux query",
            "Create a 'region' variable chained to the selected service",
            "Update panel queries to use $service and $region",
        ],
        [
            "Query variables populated from Flux live data",
            "Chained variable filters correctly on selection",
            "Dashboard reacts dynamically to dropdown changes",
        ],
        "Changing 'service' immediately narrows 'region' options — highly practical real-world pattern.",
    )

    # 7.3 Templating
    content_slide(
        "Templating: Reusable Dashboards",
        "Query  ·  Custom  ·  Interval variable types",
        [
            "Query variable — values come from real metric labels (fully dynamic)",
            "Custom variable — a fixed list defined by the dashboard author",
            "Interval variable — control the time step/resolution of queries",
            "Panel Repeat — auto-expand one panel definition into N panels per variable",
        ],
        "Show contrast: static dashboard (one service) vs templated (any service from dropdown).",
    )

    # 7.4 Lab 06
    lab_slide(
        6,
        "Templating",
        "Exercises/06-templating.md",
        [
            "Import the starter.json dashboard skeleton",
            "Create a 'service' Query variable (from Prometheus label values)",
            "Create a custom 'interval' variable: 1m, 5m, 15m, 1h",
            "Enable Panel Repeat on the main panel to expand per service",
        ],
        [
            "Query variable drives panel content dynamically",
            "Custom interval variable controls query resolution",
            "Panel Repeat expands one panel per service value",
        ],
        "Key insight: variable values come from real metric labels, not hard-coded lists.",
    )

    # 7.5 Interactive Dashboards (diagram)
    sld = prs.slides.add_slide(BLANK)
    bg(sld, W)
    header_bar(sld, "Interactive Dashboards",
               "Dashboard Links  ·  Data Links  ·  Drill-Down Flow")
    diagram_drill_down(sld)
    notes(sld, "SRE / on-call: go quickly from 'something is wrong' to 'exactly what and where'.")

    # 7.6 Lab 07
    lab_slide(
        7,
        "Interactive Dashboards",
        "Exercises/07-interactive-dashboards.md",
        [
            "Create 'Ops Overview' with Table panel: sum by (job, instance) (up)",
            "Create 'Target Detail' dashboard with time series graphs",
            "Add a Dashboard Link from Overview → Detail (preserve time range)",
            "Add a Data Link on the table row to pass $instance as variable",
        ],
        [
            "Dashboard link navigates and preserves time range",
            "Data link in Table opens correct detail view",
            "Variable ($instance) context carried between dashboards",
        ],
        "Payoff: clicking a table row lands on the detail dashboard with the correct target pre-selected.",
    )

    # ══ SECTION 8: Working with Alerting ══════════════════════════════════════
    section_divider(8, "Working with Alerting",
                    "Rules  ·  Thresholds  ·  State Transitions")

    # 8.1 Alerting System (diagram)
    sld = prs.slides.add_slide(BLANK)
    bg(sld, W)
    header_bar(sld, "Grafana Alerting System",
               "Evaluate  ·  State Machine  ·  Notify")
    diagram_alert_states(sld)
    notes(sld, "Alerting is the next step beyond dashboards — the system watches automatically instead of waiting for someone to open a graph.")

    # 8.2 Setting Up Alert Rules
    content_slide(
        "Setting Up Alert Rules",
        "Query  →  Condition  →  Pending  →  Fire",
        [
            "Write a PromQL query that measures something operationally meaningful",
            "Set a threshold condition — e.g. value > 10",
            "Evaluation interval — how often the rule is checked (e.g. every 1m)",
            "Pending period — how long condition must hold to prevent flapping",
        ],
        "Focus on the core rule evaluation cycle. No external notification channels needed in this lab.",
    )

    # 8.3 Lab 08
    lab_slide(
        8,
        "Alerting Basics",
        "Exercises/08-alerting-basics.md",
        [
            "Create an alert rule:  increase(demo_requests_total[1m]) > 10",
            "Set evaluation interval to 1m and pending period to 2m",
            "Run  ./generate-traffic.sh  to trigger the condition",
            "Watch the state transition in the Alert Rules UI",
        ],
        [
            "Alert rule created with threshold condition",
            "Traffic script successfully triggers the alert",
            "Normal → Pending → Firing transition observed in UI",
        ],
        "A PromQL expression learners already know becomes an alert with just a threshold — key bridge between observability and operations.",
    )

    # ══ SECTION 9: Grafana Plugins ════════════════════════════════════════════
    section_divider(9, "Grafana Plugins",
                    "Data Source  ·  Panel  ·  App  ·  TypeScript + React")

    # 9.1 Plugin Types (diagram)
    sld = prs.slides.add_slide(BLANK)
    bg(sld, W)
    header_bar(sld, "What are Grafana Plugins?",
               "Extend the platform for any specialized need")
    diagram_plugin_types(sld)
    notes(sld, "Shift learners from Grafana users to Grafana extenders — they now see the platform as developers.")

    # 9.2 Installing Plugins
    content_slide(
        "Installing and Configuring Plugins",
        "Public catalog  ·  Local development  ·  Version compatibility",
        [
            "Browse the public catalog at  grafana.com/plugins",
            "Install via:  grafana-cli plugins install <plugin-id>",
            "For development: mount the  dist/  folder via Docker Compose volume",
            "Verify permissions, settings, and Grafana version compatibility after install",
        ],
        "Before writing a plugin, understand the installation experience your users will have.",
    )

    # 9.3 Developing a Panel Plugin
    content_slide(
        "Developing a Basic Panel Plugin",
        "Scaffold  ·  Develop  ·  Live Reload",
        [
            "Scaffold with:   npx @grafana/create-plugin@latest",
            "Develop in TypeScript + React — same stack as Grafana core",
            "Run  npm run dev  to start the file watcher (live reload)",
            "Mount  dist/  in docker-compose.yml — no container rebuild on changes",
        ],
        "The generator provides correct project structure, Compose setup, and TypeScript config from the start.",
    )

    # 9.4 Lab 09
    lab_slide(
        9,
        "Panel Plugin Basics",
        "Exercises/09-panel-plugin-basics.md",
        [
            "Run  npx @grafana/create-plugin@latest  in labfiles/09-panel-plugin-basics/workspace/",
            "Select 'Panel' plugin type and enter a plugin name",
            "Run  npm install && npm run dev",
            "Edit  SimplePanel.tsx, save, and confirm live reload in Grafana",
        ],
        [
            "Plugin scaffolded and installed successfully",
            "Grafana detects and loads the new plugin",
            "Code edit reflected in rendered panel (live reload working)",
        ],
        "Every learner should see their own plugin render at least once — make this moment count.",
    )

    # ══ CONCLUSION ════════════════════════════════════════════════════════════
    section_divider(10, "Summary & Next Steps",
                    "Recap  ·  Q&A  ·  Keep building")

    # Summary
    content_slide_sm(
        "Course Summary",
        "What you built over two days",
        [
            "Launched Grafana in Codespaces using Docker Compose",
            "Connected Prometheus (PromQL) and InfluxDB (Flux) as data sources",
            "Built dashboards with panels, variables, and templated views",
            "Created a drill-down flow with Dashboard and Data Links",
            "Configured alert rules and observed state transitions",
            "Scaffolded and rendered a custom panel plugin",
        ],
        "Crisp recap: every topic is part of one path — from system setup to self-directed platform extension.",
    )

    # Q&A & Next Steps
    content_slide(
        "Q&A and Next Steps",
        "Keep building beyond the course",
        [
            "Swap in your team's real Prometheus / InfluxDB data sources",
            "Build dashboards that answer your team's actual operational questions",
            "Extend the panel plugin with custom options and data transformations",
            "Explore Grafana Loki (logs) and Tempo (traces) for full-stack observability",
        ],
        "Close with an open invitation — encourage learners to take what they built back to real work.",
    )

    # ── Empty Placeholder Slide ───────────────────────────────────────────────
    sld = prs.slides.add_slide(BLANK)
    bg(sld, W)
    s = rect(sld, 0.4, 0.35, 12.53, 1.25, fill=LB, line=PB, lw=Pt(1.5))
    stxt(s, "[ Slide Title — click to customise ]", size=26,
         color=RGBColor(0x99, 0xBB, 0xDD), bold=True)
    s = rect(sld, 0.4, 1.85, 12.53, 5.2,
             fill=RGBColor(0xFA, 0xFB, 0xFF),
             line=RGBColor(0xCC, 0xDD, 0xFF), lw=Pt(1.5))
    stxt(s, "[ Content area — add text, diagrams, images, or charts here ]",
         size=20, color=RGBColor(0xAA, 0xBB, 0xCC))
    rect(sld, 0.4, 7.1, 12.53, 0.3, fill=LB)
    notes(sld, "Blank template slide. Duplicate and customise as needed.")

    # ── Save ──────────────────────────────────────────────────────────────────
    out_path = os.path.join(
        os.path.dirname(os.path.abspath(__file__)),
        "grafana-for-developer-course-presentation.pptx"
    )
    prs.save(out_path)
    slide_count = len(prs.slides)
    print(f"✓ Saved {slide_count} slides → {out_path}")


if __name__ == "__main__":
    main()
